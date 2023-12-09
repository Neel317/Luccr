# import required things 
from pptx import Presentation 
from pptx.util import Inches, Pt
from python_pptx_text_replacer import TextReplacer


replacer = TextReplacer("Test LUCCR.pptx", slides='7-8', tables=True, charts=True, textframes=True)
replacer.replace_text( [('Suscepti', 'Explore') ] )
replacer.write_presentation_to_file("./Testing Luccr.pptx")


def search_and_replace(search_str, repl_str, input, output):
    """"search and replace text in PowerPoint while preserving formatting"""
    #Useful Links ;)
    #https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    #https://stackoverflow.com/questions/45247042/how-to-keep-original-text-formatting-of-text-with-python-powerpoint
    from pptx import Presentation
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                print(shape.text_frame)
                print(shape.text_frame.paragraphs[0].runs[0].text)
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(output)

def replace_text(replacements, shapes):
    """Takes dict of {match: replacement, ... } and replaces all matches.
    Currently not implemented for charts or graphics.
    """
    for shape in shapes:
        for match, replacement in replacements.items():
            if shape.has_text_frame:
                if (shape.text.find(match)) != -1:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        whole_text = "".join(run.text for run in paragraph.runs)
                        whole_text = whole_text.replace(str(match), str(replacement))
                        for idx, run in enumerate(paragraph.runs):
                            if idx != 0:
                                p = paragraph._p
                                p.remove(run._r)
                        if bool(paragraph.runs):
                            paragraph.runs[0].text = whole_text



# search_and_replace('Inside', 'Outside', './Test.pptx', './Test.pptx')
# prs = Presentation('Test.pptx')
#     # To get shapes in your slides
# slides = [slide for slide in prs.slides]
# shapes = []
# for slide in slides:
#     for shape in slide.shapes:
#         shapes.append(shape)

# replaces = {
#                     '{{PPTX}}': 'XML',
#                     '{{blah}}': 'kala'
#             }
# replace_text(replaces, shapes)
# prs.save('Test.pptx')
print('done')